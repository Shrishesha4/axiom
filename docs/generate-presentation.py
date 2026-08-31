#!/usr/bin/env python3
"""Generate editable Axiom MVP presentation (PPTX)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

DOCS = Path(__file__).resolve().parent
SCREENSHOTS = DOCS / "presentation-screenshots"

# Brand colors (teal + dark slate)
TEAL = RGBColor(13, 148, 136)
DARK = RGBColor(30, 41, 59)
MUTED = RGBColor(100, 116, 139)
WHITE = RGBColor(255, 255, 255)
LIGHT_BG = RGBColor(248, 250, 252)


def set_slide_bg(slide, color: RGBColor = LIGHT_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_bar(slide, title: str, subtitle: str | None = None):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.05))  # rectangle
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.18), Inches(9), Inches(0.55))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.5), Inches(0.72), Inches(9), Inches(0.3))
        stf = sub.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(12)
        sp.font.color.rgb = RGBColor(204, 251, 241)


def add_bullets(slide, items: list[str], left=0.55, top=1.35, width=8.9, height=5.5, size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = DARK
        p.space_after = Pt(10)
        p.bullet = True


def add_body_text(slide, text: str, left=0.55, top=1.35, width=8.9, height=1.2, size=16, bold=False):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = DARK


def add_image_if_exists(slide, filename: str, left, top, width):
    path = SCREENSHOTS / filename
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))


def build_presentation(output_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, WHITE)
    title = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(8.4), Inches(1.2))
    tp = title.text_frame.paragraphs[0]
    tp.text = "Axiom"
    tp.font.size = Pt(54)
    tp.font.bold = True
    tp.font.color.rgb = TEAL

    sub = slide.shapes.add_textbox(Inches(0.8), Inches(3.1), Inches(8.4), Inches(1.5))
    stf = sub.text_frame
    stf.word_wrap = True
    p1 = stf.paragraphs[0]
    p1.text = "AI-Powered Competitive Intelligence for Life Sciences"
    p1.font.size = Pt(24)
    p1.font.color.rgb = DARK
    p2 = stf.add_paragraph()
    p2.text = "MVP Prototype Presentation"
    p2.font.size = Pt(18)
    p2.font.color.rgb = MUTED
    p2.space_before = Pt(12)

    meta = slide.shapes.add_textbox(Inches(0.8), Inches(5.8), Inches(8.4), Inches(0.8))
    mp = meta.text_frame.paragraphs[0]
    mp.text = "Life Sciences  |  Cloud  |  AI-Assisted Development"
    mp.font.size = Pt(14)
    mp.font.color.rgb = MUTED

    # --- Slide 2: Use Case / Problem Statement ---
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title_bar(slide, "Use Case / Problem Statement")
    add_body_text(
        slide,
        "Life sciences teams need fast competitive landscape views, but data is scattered across multiple public sources.",
        height=0.8,
    )
    add_bullets(
        slide,
        [
            "Researchers manually search ClinicalTrials.gov, openFDA, and PubMed - often exporting spreadsheets",
            "Competitive intelligence takes hours and becomes stale quickly",
            "Strategy and BD teams need answers like: Who is running trials? Which mechanisms are crowded? Where is white-space?",
            "Generic AI chatbots are risky - they can invent trial counts and drug statistics",
            "There is no single workspace that fetches live public data, visualizes it, and explains it safely",
        ],
        top=2.0,
        size=17,
    )
    add_body_text(
        slide,
        "Selected use case: Competitive pharmaceutical intelligence powered by live public APIs + AI agent.",
        top=6.2,
        height=0.6,
        size=15,
        bold=True,
    )

    # --- Slide 3: Approach ---
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title_bar(slide, "Approach")
    add_bullets(
        slide,
        [
            "Build an MVP web app that accepts plain-English research questions",
            "Use an AI agent with tool calling - the LLM decides which data tools to run",
            "Fetch live data at query time from ClinicalTrials.gov, openFDA, and PubMed",
            "Compute metrics in Python (momentum scores, white-space signals) - not in the LLM",
            "Generate dynamic dashboards, follow-up chat, executive briefings, and bull/bear debate",
            "Package everything in Docker; deploy on GCP VM with nginx and a free dpdns domain",
            "Vibe-code rapidly with Cursor while keeping architecture and guardrails human-designed",
        ],
        top=1.35,
        size=17,
    )

    # --- Slide 4: Architecture ---
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title_bar(slide, "Architecture Overview")
    arch_box = slide.shapes.add_textbox(Inches(0.55), Inches(1.4), Inches(4.2), Inches(5.5))
    atf = arch_box.text_frame
    atf.word_wrap = True
    arch_lines = [
        "Frontend: Next.js + React",
        "Backend: FastAPI (Python)",
        "Database: PostgreSQL",
        "AI: OpenRouter + DeepSeek V4 Flash",
        "Data: ClinicalTrials.gov, openFDA, PubMed",
        "Deploy: Docker, GCP VM, nginx",
        "",
        "Flow:",
        "User question -> Agent -> Tools -> Live APIs",
        "-> Dashboard + AI explanation",
    ]
    for i, line in enumerate(arch_lines):
        p = atf.paragraphs[0] if i == 0 else atf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = DARK
        p.space_after = Pt(6)
        if line.endswith(":"):
            p.font.bold = True

    # Simple architecture diagram using shapes
    layers = [
        ("Browser / Next.js", 5.0, 1.6),
        ("FastAPI Agent + Tools", 5.0, 2.7),
        ("PostgreSQL", 5.0, 3.8),
        ("Public APIs + DeepSeek V4 Flash", 5.0, 4.9),
    ]
    for label, x, y in layers:
        shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(4.0), Inches(0.75))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(204, 251, 241)
        shape.line.color.rgb = TEAL
        tf = shape.text_frame
        tf.text = label
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = DARK
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # --- Slide 5: Key Features ---
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title_bar(slide, "Key Features")
    add_bullets(
        slide,
        [
            "Live investigation - fresh trial, FDA, and PubMed data per query",
            "AI agent with 8+ tools (search trials, landscape, momentum, white-space, safety, publications)",
            "Dynamic dashboard - charts adapt to question intent",
            "Server-computed momentum rankings and white-space opportunity detection",
            "Explain these signals - AI narrates dashboard data with citations",
            "Follow-up chat on the same investigation",
            "Executive briefing export and bull vs bear investment debate",
            "Session library, portfolio view, auth, and admin token limits",
        ],
        top=1.35,
        size=16,
    )

    # --- Slide 6: Demo - Home ---
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title_bar(slide, "Solution Demonstration", "Step 1 - Ask a research question")
    add_image_if_exists(slide, "01-home.png", 0.45, 1.2, 9.1)
    add_body_text(
        slide,
        "User enters a plain-English question or picks an example prompt from the command center.",
        top=6.55,
        height=0.5,
        size=14,
    )

    # --- Slide 7: Demo - Workspace ---
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title_bar(slide, "Solution Demonstration", "Step 2 - Live dashboard and AI agent")
    add_image_if_exists(slide, "02-workspace-dashboard.png", 0.35, 1.15, 9.3)
    add_body_text(
        slide,
        "Agent fetches live data, builds KPIs/charts, runs bull/bear debate, and supports follow-up questions.",
        top=6.55,
        height=0.5,
        size=14,
    )

    # --- Slide 8: Demo - Library ---
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title_bar(slide, "Solution Demonstration", "Step 3 - Saved investigations")
    add_image_if_exists(slide, "03-library.png", 0.45, 1.2, 9.1)
    add_body_text(
        slide,
        "Users revisit past investigations from the library and portfolio views.",
        top=6.55,
        height=0.5,
        size=14,
    )

    # --- Slide 9: How it addresses the problem ---
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title_bar(slide, "How the Prototype Solves the Problem")
    add_bullets(
        slide,
        [
            "One question replaces hours of manual searching across multiple websites",
            "Data is live - fetched at investigation time, not from a stale demo dataset",
            "Visual dashboard makes competitive landscape easy to scan",
            "AI explains results but numbers come from code and APIs - reducing hallucination risk",
            "Briefings and debate views support real strategy conversations",
            "Cloud-ready Docker deployment proves the app can run in production-like hosting",
        ],
        top=1.35,
        size=17,
    )

    # --- Slide 10: Key Learnings & Challenges ---
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title_bar(slide, "Key Learnings & Challenges")
    add_bullets(
        slide,
        [
            "Challenge: LLMs want to answer from memory - solved with tool-first agent design and strict prompts",
            "Challenge: Live API calls are slow (30-60s) - addressed with streaming UI and agent trace visibility",
            "Challenge: Dashboard must match question intent - built intent classification for dynamic sections",
            "Challenge: Deploying full stack on a VM - solved with Docker Compose + nginx reverse proxy",
            "Learning: Server-computed metrics are more trustworthy than LLM-generated scores",
            "Learning: MVP scope matters - public APIs first, enterprise features later",
            "Learning: Vibe coding is fast, but testing real data flows is still essential",
        ],
        top=1.35,
        size=16,
    )

    # --- Slide 11: AI-Assisted Development ---
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title_bar(slide, "AI-Assisted Development")
    add_body_text(
        slide,
        "This MVP was vibe-coded - built rapidly with AI coding tools while keeping human ownership of architecture and product decisions.",
        top=1.35,
        height=0.9,
        size=16,
    )
    add_bullets(
        slide,
        [
            "Cursor - AI pair programmer for frontend, backend, Docker, and debugging",
            "DeepSeek V4 Flash (open-source via OpenRouter) - agent reasoning, tool selection, summaries, briefings, debate",
            "OpenRouter - one API for model access, easy swapping, optional web search plugins",
            "AI helped write boilerplate UI, API routes, and agent orchestration code",
            "Human decisions: data sources, tool schemas, prompts, guardrails, deployment, and UX flow",
            "Validation: real NCT IDs, matching chart data, health checks, and end-to-end investigation tests",
        ],
        top=2.2,
        size=16,
    )

    # --- Slide 12: Roadmap ---
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    add_title_bar(slide, "Next Steps / Roadmap")
    add_bullets(
        slide,
        [
            "Phase 1 (Done): MVP with live public data, agent, dashboard, briefing, GCP VM deploy",
            "Phase 2: Caching, monitoring, CI/CD, more data sources, mobile polish",
            "Phase 3: Enterprise - private data connectors, alerts, SSO, audit logs, multi-tenant GCP",
            "Future: RAG on internal documents, regulatory-ready exports, forecasting module",
        ],
        top=1.35,
        size=18,
    )

    # --- Slide 13: Thank you ---
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, TEAL)
    thanks = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(8.4), Inches(1.2))
    tp = thanks.text_frame.paragraphs[0]
    tp.text = "Thank You"
    tp.font.size = Pt(48)
    tp.font.bold = True
    tp.font.color.rgb = WHITE
    tp.alignment = PP_ALIGN.CENTER

    qa = slide.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(8.4), Inches(0.8))
    qp = qa.text_frame.paragraphs[0]
    qp.text = "Questions & Demo"
    qp.font.size = Pt(22)
    qp.font.color.rgb = RGBColor(204, 251, 241)
    qp.alignment = PP_ALIGN.CENTER

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


if __name__ == "__main__":
    out = DOCS / "Axiom-MVP-Presentation.pptx"
    build_presentation(out)
    print(f"Created: {out}")
